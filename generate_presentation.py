from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import datetime

# Create presentation
prs = Presentation()
prs.core_properties.title = "Sari-Sari POS — Project Defense"
prs.core_properties.author = "Sari-Sari POS Team"

# Colors (matching site): dark navy, leaf green, soft blue, white
DARK = RGBColor(0x0F, 0x17, 0x2A)   # #0f172a
LEAF = RGBColor(0x4A, 0x7C, 0x59)   # #4a7c59
SOFT = RGBColor(0xC3, 0xDA, 0xFE)   # #c3dafe
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x06, 0xB6, 0xD4) # #06b6d4

# Helpers
def add_title_slide(title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_tf = slide.shapes.title.text_frame
    subtitle_tf = slide.placeholders[1].text_frame
    title_tf.text = title
    p = subtitle_tf.paragraphs[0]
    p.text = subtitle
    # color title
    title_tf.paragraphs[0].runs[0].font.size = Pt(40)
    title_tf.paragraphs[0].runs[0].font.color.rgb = DARK
    subtitle_tf.paragraphs[0].runs[0].font.color.rgb = LEAF
    return slide

def add_bullet_slide(title, bullets, notes=None):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
            p.text = b
        else:
            p = body.add_paragraph()
            p.text = b
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = DARK
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
    return slide

# Title
today = datetime.date.today().strftime('%B %d, %Y')
add_title_slide('Sari-Sari POS', f'Project Defense — {today}\nTeam: Maquimot & Nu\u00f1ez')
prs.slides[-1].notes_slide.notes_text_frame.text = 'Introduce yourselves, quick one-liner about the system\nMention live demo at the end.'

# Agenda
add_bullet_slide('Agenda', [
    'Motivation & Problem Statement',
    'Solution Overview',
    'Key Features & Workflow',
    'Architecture & Data Model',
    'Demo walkthrough',
    'Testing, Security, Future work',
    'Q&A'
], notes='Walk through each agenda item briefly. Keep time for demo and Q&A.')

# Problem
add_bullet_slide('Problem Statement', [
    'Small stores need a lightweight POS to manage inventory and sales',
    'Manual record-keeping causes errors and wastes time',
    'Need offline-friendly, simple UI for quick transactions'
], notes='Give real-world context - sari-sari store pain points.')

# Solution Overview
add_bullet_slide('Solution Overview', [
    'Web-based POS built with Django and SQLite',
    'Quick product CRUD, sales recording, and daily sales reporting',
    'Simple authentication and role-based access for admin actions'
], notes='Highlight how the solution addresses the problems listed.')

# Key Features
add_bullet_slide('Key Features', [
    'Secure login and dashboard (protected routes)',
    'Inventory management: add/edit/restock products and suppliers',
    'Record sales with multiple payment methods (Cash, GCash, Utang)',
    'Low-stock warnings and daily sales reports',
    'Debt (utang) tracking and marking as paid via dashboard'
], notes='For each feature, be ready with a 15–30s demo snippet or screenshot.')

# Architecture
add_bullet_slide('Architecture', [
    'Django (views, models, templates) as backend and server-rendered UI',
    'SQLite for local persistent storage (db.sqlite3)',
    'Bootstrap + FontAwesome for responsive UI',
    'Simple URL routing and auth via django.contrib.auth'
], notes='Show file layout and where main modules live (POS app).')

# Data Model
add_bullet_slide('Data Model (Main models)', [
    'Product (name, price, stock_quantity, category, supplier)',
    'Category',
    'Supplier',
    'Transaction (payment_method, total_amount, is_paid, timestamp)',
    'TransactionItem (transaction, product_name, price, quantity)'
], notes='Explain relationships: Transaction -> TransactionItem (1:N).')

# UI/Color Scheme Visuals
slide_layout = prs.slide_layouts[6]  # blank
slide = prs.slides.add_slide(slide_layout)
# Add title
tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7)).text_frame
tx.text = 'UI & Color Scheme'
tx.paragraphs[0].runs[0].font.size = Pt(28)
tx.paragraphs[0].runs[0].font.color.rgb = DARK

# add color swatches
left = Inches(0.6)
for i, (name, rgb) in enumerate([('Dark Navy', DARK), ('Leaf Green', LEAF), ('Soft Blue', SOFT), ('Accent', ACCENT)]):
    shape = slide.shapes.add_shape(1, left + Inches(i*1.6), Inches(1.2), Inches(1.2), Inches(1.2))
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = rgb
    # label
    lbl = slide.shapes.add_textbox(left + Inches(i*1.6), Inches(2.5), Inches(1.6), Inches(0.4)).text_frame
    lbl.text = name
    lbl.paragraphs[0].runs[0].font.size = Pt(12)
    lbl.paragraphs[0].runs[0].font.color.rgb = DARK

slide.notes_slide.notes_text_frame.text = 'Match this palette during live demo. Show login and dashboard colors.'

# Demo Steps
add_bullet_slide('Demo Walkthrough', [
    '1) Show login flow (login screen first)',
    '2) Open dashboard: list of products and low-stock warnings',
    '3) Record a sale (Cash and Utang examples)',
    '4) Show daily sales report and mark utang as paid',
    '5) Manage inventory: add / edit / restock product'
], notes='Perform live demo in this sequence. Mention expected results and where code lives.')

# Security & Testing
add_bullet_slide('Security & Testing', [
    'Authentication uses Django auth (login required decorators)',
    'Input validation and atomic transactions for sales',
    'Unit tests (if present) and manual QA flows',
    'Backups: db.sqlite3 copy strategy for deployments'
], notes='Be prepared to answer questions about auth and data integrity.')

# Performance & Scalability
add_bullet_slide('Performance & Scalability', [
    'Designed for small stores; SQLite gives simplicity for local deploy',
    'For scale: swap to PostgreSQL, add caching, and deploy behind WSGI server',
    'Separate services or microservices not necessary for MVP'
], notes='Have migration path ready if asked about large-usage scenarios.')

# Future Work / Roadmap
add_bullet_slide('Roadmap & Future Work', [
    'Add user roles and permissions (cashiers vs managers)',
    'Export reports (CSV/PDF), scheduled backups, and analytics',
    'Mobile-friendly checkout and offline sync',
    'Integrate with simple payment APIs (GCash receipt verification)'
], notes='Prioritize features by business value for the panel.' )

# Backup: Key Code Paths
add_bullet_slide('Key Code Paths (Pointers)', [
    'POS/views.py — dashboard, inventory, record_sale, utang API',
    'POS/models.py — product and transaction models',
    'templates/ — UI: POS/base.html, POS/login.html, POS/dashboard.html',
    'posSystem/settings.py — LOGIN_REDIRECT_URL and static settings'
], notes='Open these files quickly during Q&A if asked about implementation details.')

# Q&A
add_bullet_slide('Q & A', ['Thank you! Questions and feedback.'], notes='Open floor for panelist questions. Have demo ready.')

# Save file
out_path = r"C:\Users\USER\OneDrive\Desktop\pos_workspace\presentation_SariSari_POS.pptx"
prs.save(out_path)
print('SAVED:', out_path)
